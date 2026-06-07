"""
BeyondSelf — Full Pitch Deck Generator
Matches the visual style from the provided screenshots exactly.
Run: python generate_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Patch RGBColor to expose .red .green .blue (missing in some versions)
RGBColor.red   = property(lambda self: self[0])
RGBColor.green = property(lambda self: self[1])
RGBColor.blue  = property(lambda self: self[2])

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height
BLANK = prs.slide_layouts[6]

# ── Palette ──────────────────────────────────────────────────────────────────
BG      = RGBColor(0x09,0x08,0x1E)
CARD    = RGBColor(0x0F,0x0E,0x2A)
CARD2   = RGBColor(0x14,0x13,0x32)
DIV     = RGBColor(0x24,0x22,0x42)
PURPLE  = RGBColor(0x8B,0x5C,0xF6)
PURPL   = RGBColor(0xC4,0xB5,0xFD)
INDIGO  = RGBColor(0x6D,0x28,0xD9)
CYAN    = RGBColor(0x22,0xD3,0xEE)
AMBER   = RGBColor(0xF5,0x9E,0x0B)
GREEN   = RGBColor(0x4A,0xDE,0x80)
ORANGE  = RGBColor(0xF9,0x71,0x16)
PINK    = RGBColor(0xEC,0x48,0x99)
BLUE    = RGBColor(0x60,0xA5,0xFA)
RED     = RGBColor(0xEF,0x44,0x44)
TEAL    = RGBColor(0x10,0xB9,0x81)
YELLOW  = RGBColor(0xFB,0xBF,0x24)
WHITE   = RGBColor(0xFF,0xFF,0xFF)
LTXT    = RGBColor(0xE2,0xE8,0xF0)
MTXT    = RGBColor(0x94,0xA3,0xB8)
DTXT    = RGBColor(0x4B,0x55,0x63)

def _i(x): return int(x)

# ── Core drawing helpers ──────────────────────────────────────────────────────

def set_bg(slide, color=BG):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def R(slide, l,t,w,h, fill=None, line=None, lw=Pt(1)):
    sh = slide.shapes.add_shape(1, _i(l),_i(t),_i(w),_i(h))
    sh.fill.solid() if fill else sh.fill.background()
    if fill: sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line if line else RGBColor(0,0,0)
    sh.line.fill.background() if not line else None
    if line: sh.line.color.rgb = line; sh.line.width = _i(lw)
    sh.shadow.inherit = False
    return sh

def RR(slide, l,t,w,h, fill=None, line=None, lw=Pt(1), r=0.06):
    sh = slide.shapes.add_shape(5, _i(l),_i(t),_i(w),_i(h))
    sh.adjustments[0] = r
    sh.fill.solid() if fill else sh.fill.background()
    if fill: sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    if line: sh.line.color.rgb = line; sh.line.width = _i(lw)
    sh.shadow.inherit = False
    return sh

def OV(slide, l,t,w,h, fill=None, line=None, lw=Pt(2)):
    sh = slide.shapes.add_shape(9, _i(l),_i(t),_i(w),_i(h))
    sh.fill.solid() if fill else sh.fill.background()
    if fill: sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    if line: sh.line.color.rgb = line; sh.line.width = _i(lw)
    sh.shadow.inherit = False
    return sh

def T(slide, text, l,t,w,h, sz=11, bold=False, color=WHITE,
      align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(_i(l),_i(t),_i(w),_i(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    rn = p.add_run(); rn.text = text
    rn.font.size = Pt(sz); rn.font.bold = bold
    rn.font.italic = italic; rn.font.color.rgb = color
    rn.font.name = 'Calibri'
    return tb

def Tlines(slide, lines, l,t,w,h):
    """lines = list of (text,sz,bold,color,align)"""
    tb = slide.shapes.add_textbox(_i(l),_i(t),_i(w),_i(h))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in lines:
        txt,sz,bold,col,al = item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = al
        rn = p.add_run(); rn.text = txt
        rn.font.size = Pt(sz); rn.font.bold = bold
        rn.font.color.rgb = col; rn.font.name = 'Calibri'

def BAR(slide,l,t,w, color, h=Pt(2)):
    R(slide, l,t,w,_i(h), fill=color)

def TAG(slide, text, l, t, bg, fg, sz=7.5, pad_x=Pt(7), h=Pt(16)):
    w = len(text)*Pt(5.5)+pad_x*2
    RR(slide, l,t,w,_i(h), fill=bg, r=0.5)
    T(slide, text, _i(l)+_i(Pt(4)),_i(t)+_i(Pt(1.5)),
      _i(w)-_i(Pt(8)),_i(h), sz=sz, bold=True, color=fg, align=PP_ALIGN.CENTER)
    return _i(l)+_i(w)+_i(Pt(5))  # next x

def slide_header(slide, eye, title, color=PURPLE, subtitle=None, brand="BeyondSelf · Digital Twin"):
    T(slide, eye, Inches(0.5), Inches(0.18), W-Inches(3.5), Inches(0.28),
      sz=8, bold=True, color=RGBColor(min(255,color.red//2+80),
                                       min(255,color.green//2+80),
                                       min(255,color.blue//2+80)))
    T(slide, title, Inches(0.5), Inches(0.44), W-Inches(3.5), Inches(0.75),
      sz=30, bold=True, color=WHITE)
    if subtitle:
        T(slide, subtitle, Inches(0.5), Inches(1.15), W-Inches(1), Inches(0.28),
          sz=11, color=MTXT)
    BAR(slide, Inches(0.5), Inches(1.45), W-Inches(1), DIV)
    # Brand pill top-right
    bw = Inches(2.6)
    RR(slide, W-bw-Inches(0.35), Inches(0.2), bw, Inches(0.35),
       fill=RGBColor(0x1A,0x18,0x38), line=PURPLE, lw=Pt(1), r=0.4)
    T(slide, brand, W-bw-Inches(0.25), Inches(0.22), bw-Inches(0.2), Inches(0.3),
      sz=8, bold=True, color=PURPL, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

# Decorative orbs
for (ox,oy,ow,oc) in [
    (Inches(7.5),Inches(-0.8),Inches(5), RGBColor(0x14,0x10,0x35)),
    (Inches(9),  Inches(4.5), Inches(4.5),RGBColor(0x10,0x0D,0x2E)),
]:
    OV(s,ox,oy,ow,ow, fill=oc)

# Left accent strip
R(s, 0,0, Inches(0.06),H, fill=PURPLE)
R(s, Inches(0.06),0, Inches(0.025),H, fill=INDIGO)

# BEYONDSELF
T(s,"BEYONDSELF", Inches(0.7),Inches(1.2), Inches(8),Inches(1.4),
  sz=64,bold=True,color=PURPLE)

# Subtitle
T(s,"Your AI-Powered Personal Digital Twin",
  Inches(0.7),Inches(2.65),Inches(9),Inches(0.55),
  sz=22,color=PURPL,italic=True)

# Tagline
T(s,"ONE TWIN.  LIMITLESS POSSIBILITIES.",
  Inches(0.7),Inches(3.3),Inches(9),Inches(0.4),
  sz=16,bold=True,color=WHITE)

BAR(s, Inches(0.7),Inches(3.85),Inches(5.5),PURPLE, h=Pt(2))

# Date + team
T(s,"4 June 2026", Inches(0.7),Inches(4.0),Inches(5),Inches(0.3), sz=12,color=MTXT)
T(s,"Presented by", Inches(0.7),Inches(4.4),Inches(4),Inches(0.28), sz=9,color=DTXT)
T(s,"Khushi Bansal  ·  Pavani Gubba  ·  Samridhi Pandey",
  Inches(0.7),Inches(4.68),Inches(9),Inches(0.35), sz=14,bold=True,color=PURPL)

# Domain badges bottom-left
for label,col,bg in [
    ("🏥 Health",  GREEN,  RGBColor(0x06,0x2A,0x14)),
    ("💰 Finance", AMBER,  RGBColor(0x33,0x1A,0x03)),
    ("💼 Career",  BLUE,   RGBColor(0x0E,0x24,0x50)),
]:
    RR(s, Inches(0.7 if label.startswith("🏥") else
                  2.5 if label.startswith("💰") else 4.3),
       Inches(5.3),Inches(1.55),Inches(0.38),
       fill=bg, line=col, lw=Pt(1.5), r=0.3)
    T(s, label,
      Inches(0.78 if label.startswith("🏥") else
             2.58 if label.startswith("💰") else 4.38),
      Inches(5.32),Inches(1.4),Inches(0.34),
      sz=11,bold=True,color=col,align=PP_ALIGN.CENTER)

# Right visual — human figure rings
for sz_in, alpha in [(3.2,0x14),(2.4,0x22),(1.6,0x35)]:
    cx = Inches(10.6 - sz_in/2); cy = Inches(4.1 - sz_in/2)
    OV(s, cx,cy,Inches(sz_in),Inches(sz_in),
       fill=None, line=RGBColor(0x6D,0x28,alpha), lw=Pt(1.2))

T(s,"◎",Inches(9.6),Inches(2.0),Inches(2.3),Inches(2.3),
  sz=88,color=RGBColor(0x4F,0x46,0xE5),align=PP_ALIGN.CENTER)
T(s,"Digital Twin",Inches(9.5),Inches(3.9),Inches(2.5),Inches(0.38),
  sz=10,color=DTXT,align=PP_ALIGN.CENTER)

# Domain icons around figure
for (di,dc,dxi,dyi) in [
    ("♥",  RGBColor(0xEF,0x44,0x44), Inches(7.4),  Inches(2.5)),
    ("💼", BLUE,                      Inches(11.8), Inches(2.5)),
    ("💳", AMBER,                     Inches(8.5),  Inches(4.8)),
    ("🧠", TEAL,                      Inches(11.0), Inches(4.8)),
]:
    RR(s, dxi,dyi,Inches(0.85),Inches(0.85),
       fill=RGBColor(dc.red//4,dc.green//4,dc.blue//4), line=dc, lw=Pt(1.5), r=0.25)
    T(s, di, dxi,dyi,Inches(0.85),Inches(0.85), sz=22, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)
slide_header(s,"THE CHALLENGE  ·  WHY BEYONDSELF EXISTS",
             "Life Data Lives in Silos.", AMBER,
             "No existing app connects Health, Finance & Career into one intelligent system.")

problems = [
    ("01","🔀","Data Fragmentation",
     "Fitbit, bank slips, LinkedIn, lab reports — 6+ disconnected apps. No single source of truth about your life.",
     AMBER, RGBColor(0x3A,0x20,0x06)),
    ("02","🕸️","Cross-Domain Blindness",
     "Higher salary → less sleep → poor health → lower output. No app modelled this invisible chain.",
     ORANGE, RGBColor(0x38,0x16,0x06)),
    ("03","🎙️","Smart Input Classification",
     "'Spent ₹450 on medicine' — Health? Finance? Both? Zero apps auto-classified natural speech.",
     PURPLE, RGBColor(0x25,0x15,0x48)),
    ("04","🔓","Privacy Exposure",
     "Raw bank statements & medical reports sent unfiltered to AI models. Names, phones, account numbers exposed.",
     RED, RGBColor(0x40,0x0A,0x10)),
]

for i,(num,icon,title,desc,col,bg) in enumerate(problems):
    cx = i%2; cy = i//2
    cl = Inches(0.45 + cx*6.52)
    ct = Inches(1.55 + cy*2.68)
    cw = Inches(6.2); ch = Inches(2.45)
    RR(s, cl,ct,cw,ch, fill=bg, line=col, lw=Pt(1.2), r=0.04)
    BAR(s, cl,ct,cw,col, h=Pt(2.5))
    T(s, num, cl+Inches(0.18),ct+Inches(0.12),Inches(0.55),Inches(0.55),
      sz=24,bold=True,color=RGBColor(col.red//4,col.green//4,col.blue//4))
    T(s, f"{icon}  {title}", cl+Inches(0.75),ct+Inches(0.14),cw-Inches(0.95),Inches(0.42),
      sz=15,bold=True,color=col)
    T(s, desc, cl+Inches(0.2),ct+Inches(0.6),cw-Inches(0.4),Inches(1.6),
      sz=11,color=MTXT,wrap=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)
slide_header(s,"THE SOLUTION  ·  PERSONAL DIGITAL TWIN",
             "BeyondSelf — One Intelligent System.", PURPLE,
             "AI-powered, privacy-first, cross-domain life intelligence platform.")

# Central twin box
RR(s, Inches(5.4),Inches(1.6),Inches(2.8),Inches(3.5),
   fill=RGBColor(0x18,0x15,0x38), line=PURPLE, lw=Pt(2), r=0.06)
BAR(s, Inches(5.4),Inches(1.6),Inches(2.8),PURPLE)
T(s,"◎", Inches(5.4),Inches(1.7),Inches(2.8),Inches(1.2),
  sz=46,color=PURPLE,align=PP_ALIGN.CENTER)
T(s,"Digital Twin", Inches(5.4),Inches(2.88),Inches(2.8),Inches(0.38),
  sz=14,bold=True,color=PURPL,align=PP_ALIGN.CENTER)
T(s,"Life Score: 74 / 100", Inches(5.4),Inches(3.28),Inches(2.8),Inches(0.32),
  sz=11,color=CYAN,align=PP_ALIGN.CENTER)
T(s,"Powered by 9 Engines", Inches(5.4),Inches(3.62),Inches(2.8),Inches(0.28),
  sz=9,color=MTXT,align=PP_ALIGN.CENTER)

# Domain cards
for (icon,name,tags,col,dl,dt) in [
    ("♥","Health","Sleep · Workouts\nStress · Nutrition · BMI",GREEN,Inches(0.4),Inches(2.4)),
    ("💰","Finance","Salary · Expenses\nSavings · Bills · Goals",AMBER,Inches(5.4),Inches(1.55)-Inches(1.3)),
    ("💼","Career","Skills · GitHub\nCertificates · Jobs",BLUE,Inches(10.45),Inches(2.4)),
]:
    bg = RGBColor(col.red//6,col.green//6,col.blue//6)
    RR(s,dl,dt,Inches(2.4),Inches(2.0), fill=bg,line=col,lw=Pt(1.5),r=0.06)
    BAR(s,dl,dt,Inches(2.4),col)
    T(s,f"{icon}  {name}", dl+Inches(0.15),dt+Inches(0.1),Inches(2.1),Inches(0.4),
      sz=15,bold=True,color=col)
    T(s,tags, dl+Inches(0.15),dt+Inches(0.55),Inches(2.1),Inches(1.2),
      sz=10,color=MTXT)

# Connector lines
R(s, Inches(2.8),Inches(3.35),Inches(2.6),Pt(1.5), fill=RGBColor(0x2A,0x45,0x2A))
R(s, Inches(7.2),Inches(3.35),Inches(3.25),Pt(1.5), fill=RGBColor(0x2A,0x30,0x50))
R(s, Inches(6.8),Inches(1.58),Pt(1.5),Inches(0.02), fill=RGBColor(0x3A,0x28,0x08))

# Key differentiators (right col)
diffs = [
    ("🎯","One life score — Health+Finance+Career"),
    ("🤖","Groq AI classifies every input auto"),
    ("📊","12-month trajectory with confidence decay"),
    ("🔒","PII stripped before any AI call"),
    ("⚡","Cascade Map: decisions ripple cross-domain"),
]
for j,(ic,dt_) in enumerate(diffs):
    RR(s, Inches(9.0),Inches(1.58+j*1.08),Inches(4.0),Inches(0.9),
       fill=RGBColor(0x12,0x11,0x2C), line=DIV, lw=Pt(1), r=0.04)
    T(s, ic, Inches(9.1),Inches(1.65+j*1.08),Inches(0.45),Inches(0.45), sz=16)
    T(s, dt_, Inches(9.6),Inches(1.67+j*1.08),Inches(3.2),Inches(0.5),
      sz=10.5,color=MTXT,wrap=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — KEY FEATURES  (3×3 bento)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)
slide_header(s,"FEATURES  ·  ORIGINAL CAPABILITIES",
             "Features That Didn't Exist Before.", PURPLE,
             "9 breakthrough features built from scratch — no rival has this combination.")

feats = [
    ("⚖️","Life Balance Score","Health+Finance+Career → one deterministic score",PURPLE),
    ("🌌","Talk With Future Self","Groq AI answers as your projected 2028 self",RGBColor(0x60,0x5C,0xF6)),
    ("🗺️","Cascade Map","Cross-domain cause-effect modelled in real time",BLUE),
    ("⚡","Stress Test Simulator","Simulate job loss / medical crisis — see cross-domain impact",AMBER),
    ("🌱","Life Bloom Plant","Goals = plant growth. Burnout = wilting SVG",GREEN),
    ("🪞","Adaptive Digital Avatar","Avatar state shifts with your Life Score",PINK),
    ("🎙️","Universal AI Input","Speak or upload — Groq routes to correct domain",TEAL),
    ("👻","Ghost Timeline","Past ghosts + future projections on one axis",CYAN),
    ("🏪","LIFE Market","Earn purchases by completing life-goal commitments",YELLOW),
]
COLS=3
for i,(ic,nm,ds,col) in enumerate(feats):
    cx=i%COLS; cy=i//COLS
    fl=Inches(0.42+cx*4.32); ft=Inches(1.58+cy*1.85)
    fw=Inches(4.0); fh=Inches(1.7)
    bg=RGBColor(max(0,col.red//7),max(0,col.green//7),max(0,col.blue//7))
    RR(s,fl,ft,fw,fh, fill=bg,line=col,lw=Pt(1.2),r=0.04)
    BAR(s,fl,ft,fw,col,h=Pt(2.5))
    T(s,f"{ic}  {nm}", fl+Inches(0.18),ft+Inches(0.16),fw-Inches(0.36),Inches(0.4),
      sz=12,bold=True,color=col)
    T(s,ds, fl+Inches(0.18),ft+Inches(0.62),fw-Inches(0.36),Inches(0.9),
      sz=10,color=MTXT,wrap=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — TECH STACK: FRONTEND + BACKEND  (matches screenshot 2)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)
slide_header(s,"TECH STACK SELECTION & REASON · SLIDE 1 OF 2",
             "Frontend & Backend", ORANGE,
             "React 19 · Vite 8 · Tailwind v4                Spring Boot 3.2 · Java 21 · PostgreSQL",
             "BeyondSelf · Digital Twin")

def tech_panel(slide, l, t, w, h, icon_txt, header, color, sub, items):
    """items = [(name, version_tag, category_tag, desc), ...]"""
    bg = RGBColor(max(0,color.red//6),max(0,color.green//6),max(0,color.blue//6))
    RR(slide, l,t,w,h, fill=CARD, line=DIV, lw=Pt(1), r=0.04)
    BAR(slide, l,t,w, color, h=Pt(2))

    # Icon box
    ib = Inches(0.52)
    RR(slide, l+Inches(0.2),t+Inches(0.18),ib,ib,
       fill=RGBColor(max(0,color.red//5),max(0,color.green//5),max(0,color.blue//5)),
       line=color,lw=Pt(1.5),r=0.2)
    T(slide, icon_txt, l+Inches(0.2),t+Inches(0.16),ib,ib,
      sz=22,align=PP_ALIGN.CENTER,color=color)

    T(slide, header, l+Inches(0.85),t+Inches(0.2),w-Inches(1.2),Inches(0.42),
      sz=18,bold=True,color=color)
    T(slide, sub, l+w-Inches(3.2),t+Inches(0.22),Inches(3.0),Inches(0.28),
      sz=8,color=DTXT,align=PP_ALIGN.RIGHT)

    BAR(slide, l+Inches(0.15),t+Inches(0.78),w-Inches(0.3),DIV)

    for j,(name,ver,cat,desc) in enumerate(items):
        iy = t+Inches(0.95+j*1.05)
        T(slide, name, l+Inches(0.2),iy,Inches(1.6),Inches(0.35),
          sz=13,bold=True,color=WHITE)
        # version badge
        nxt = TAG(slide, ver, l+Inches(0.2),iy+Inches(0.36),
                  RGBColor(0x16,0x3B,0x1A), GREEN, sz=7)
        # category badge
        TAG(slide, cat, nxt,iy+Inches(0.36),
            RGBColor(0x2A,0x20,0x50), PURPL, sz=7)
        T(slide, desc, l+Inches(2.0),iy,w-Inches(2.3),Inches(0.85),
          sz=9.5,color=MTXT,wrap=True)

PW = Inches(6.25)
PH = Inches(5.75)
PT = Inches(1.6)

tech_panel(s, Inches(0.38),PT, PW,PH,
           "⚡","Frontend", ORANGE,
           "React 19 · Vite 8 · Tailwind v4", [
    ("React 19",      "v19.2",  "Concurrent Rendering",
     "New concurrent model keeps UI responsive while 9 scoring engines compute simultaneously in the background — no jank on the Dashboard"),
    ("Vite 8",        "v8.0",   "Build Tool",
     "Native ES modules give sub-100ms HMR; lazy-loaded heavy pages (Coach, NeuralCore) only bundle when visited"),
    ("Tailwind CSS v4","v4.3",  "Styling",
     "Utility-first lets us build the glassmorphism dark UI and custom glass-card / bg-mesh classes without writing raw CSS files"),
    ("Framer Motion", "v12.38", "Animation",
     "Declarative spring physics for score rings, page transitions, and cascade visualizations — no manual requestAnimationFrame"),
    ("Recharts",      "v3.8",   "Data Viz",
     "React-native SVG charts — AreaChart, BarcChart, RadialBarChart render domain trends with zero canvas complexity"),
])

tech_panel(s, Inches(6.7),PT, PW,PH,
           "🖥","Backend", BLUE,
           "Spring Boot 3.2 · Java 21 · PostgreSQL", [
    ("Spring Boot 3.2","Java 21","API Server",
     "Virtual threads (Project Loom) handle concurrent AI proxy calls without thread-pool exhaustion — critical for the Groq + Gemini proxy endpoints"),
    ("JWT + BCrypt",   "JWT 0.12","Auth",
     "Stateless auth supports two token formats simultaneously — real JWT for backend users and base64 demo tokens for offline judges"),
    ("Bucket4j",       "v8.10","Rate Limiting",
     "Per-IP token buckets cap AI proxy at 20 req/min — protects free-tier Groq quota from accidental or malicious exhaustion"),
    ("Flyway",         "v10.10","DB Migrations",
     "Versioned SQL scripts prevent schema drift between teammates' local databases and production — no 'works on my machine' DB issues"),
    ("POI · PDFBox · OpenCSV","—","File Parsing",
     "Server-side parse of Excel, PDF & CSV uploads with AES-256-GCM encryption at rest — medical and financial files never stored in plaintext"),
])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — AI LAYER & INTEGRATIONS  (matches screenshot 3)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)
slide_header(s,"TECH STACK SELECTION & REASON · SLIDE 2 OF 2",
             "AI Layer & Integrations", PINK,
             "Free-tier AI with deterministic fallbacks — fully functional with zero API keys",
             "BeyondSelf · Digital Twin")

def ai_panel(slide, l, t, w, h, icon_txt, header, color, sub, items):
    """items = [(abbrev, name, tag, tag_color, desc), ...]"""
    RR(slide, l,t,w,h, fill=CARD, line=DIV, lw=Pt(1), r=0.04)
    BAR(slide, l,t,w, color, h=Pt(2))
    ib = Inches(0.52)
    RR(slide, l+Inches(0.2),t+Inches(0.18),ib,ib,
       fill=RGBColor(max(0,color.red//5),max(0,color.green//5),max(0,color.blue//5)),
       line=color,lw=Pt(1.5),r=0.2)
    T(slide,icon_txt,l+Inches(0.2),t+Inches(0.16),ib,ib,sz=22,align=PP_ALIGN.CENTER,color=color)
    T(slide,header, l+Inches(0.85),t+Inches(0.2),w-Inches(1.2),Inches(0.42),
      sz=18,bold=True,color=color)
    T(slide,sub, l+w-Inches(3.4),t+Inches(0.24),Inches(3.2),Inches(0.28),
      sz=8,color=DTXT,align=PP_ALIGN.RIGHT)
    BAR(slide,l+Inches(0.15),t+Inches(0.78),w-Inches(0.3),DIV)

    for j,(abbr,name,tag,tc,desc) in enumerate(items):
        iy = t+Inches(0.95+j*1.05)
        # Abbr badge
        RR(slide, l+Inches(0.2),iy+Inches(0.03),Inches(0.48),Inches(0.42),
           fill=RGBColor(tc.red//5,tc.green//5,tc.blue//5),line=tc,lw=Pt(1),r=0.2)
        T(slide,abbr, l+Inches(0.2),iy,Inches(0.48),Inches(0.48),
          sz=8,bold=True,color=tc,align=PP_ALIGN.CENTER)
        T(slide,name, l+Inches(0.82),iy,Inches(1.9),Inches(0.35),sz=13,bold=True,color=WHITE)
        TAG(slide,tag, l+Inches(0.82),iy+Inches(0.37),
            RGBColor(tc.red//5,tc.green//5,tc.blue//5),tc,sz=7)
        T(slide,desc, l+Inches(2.1),iy,w-Inches(2.35),Inches(0.85),sz=9.5,color=MTXT,wrap=True)

ai_panel(s, Inches(0.38),Inches(1.6), PW,PH,
         "🧠","AI & Intelligence", PINK,
         "Groq · Gemini · Tesseract", [
    ("GRQ","Groq · LLaMA 3.3-70b","Chat / Coach",  ORANGE,
     "200 tok/s — instant coach replies; free tier, same quality as OpenAI GPT-4"),
    ("L4V","LLaMA 4 Scout Vision", "Multimodal",   PURPLE,
     "Reads meal photos, medical reports & bank statements → auto-logs to correct domain"),
    ("GEM","Gemini 2.0 Flash",     "Backend Proxy", BLUE,
     "Routed through Spring Boot — key never hits the browser; PII stripped before LLM"),
    ("OCR","Tesseract.js",         "OCR",           CYAN,
     "In-browser WASM OCR for receipts & scanned PDFs — no server upload needed"),
    ("9×", "Deterministic Engines","Resilience",    GREEN,
     "Pure rule-based math — fully functional offline, zero hallucination on scores"),
])

ai_panel(s, Inches(6.7),Inches(1.6), PW,PH,
         "🔗","Integrations", AMBER,
         "GitHub · Fitbit · Jobs · Nutrition", [
    ("GH", "GitHub API v3",    "Developer",  RGBColor(0xD1,0xD5,0xDB),
     "Real contribution heatmap & language stats → auto-feeds Career Score, no manual input"),
    ("FIT","Fitbit OAuth 2.0", "Wearable",   GREEN,
     "Syncs sleep, heart rate & steps into Health Score — eliminates manual logging friction"),
    ("JOB","Adzuna + Jooble",  "Job Boards", BLUE,
     "Dual-API proxy for Indian + global listings matched to skill gaps; backend caches quota"),
    ("NX", "Nutritionix API",  "Nutrition",  ORANGE,
     "500k+ foods with accurate macros; Groq fallback when quota exceeded"),
    ("DB", "PostgreSQL",       "Persistence",PURPLE,
     "Production-grade relational DB with Flyway versioned schema and JPA entities"),
])

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — SYSTEM ARCHITECTURE  (matches screenshot 4)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)
T(s,"BeyondSelf — System Architecture", Inches(0.4),Inches(0.18),Inches(9),Inches(0.55),
  sz=24,bold=True,color=WHITE)
T(s,"Offline-first.  Deterministic AI.  Privacy by design.",
  Inches(0.4),Inches(0.72),Inches(9),Inches(0.28),sz=11,color=MTXT)
BAR(s,Inches(0),Inches(1.08),W,DIV)

# ── External APIs top bar ────────────────────────────────────────────────────
R(s,Inches(0.3),Inches(1.15),W-Inches(0.6),Inches(1.1),
  fill=RGBColor(0x12,0x10,0x2A),line=DIV,lw=Pt(1))
T(s,"EXTERNAL APIs & THIRD-PARTY SERVICES",
  Inches(0),Inches(1.2),W,Inches(0.25),sz=8,bold=True,
  color=PURPL,align=PP_ALIGN.CENTER)
apis=[
    ("Google\nGemini 2.0","Vision·NLP·Long","Server-Side Only",RGBColor(0xEA,0x43,0x35)),
    ("Groq\nLlama 3.3-70b","Voice·Nutrition·Career","Fast Inference",ORANGE),
    ("GitHub\nAPI v3","Repos·Contribs·Stars","Read Only",WHITE),
    ("NewsData\nAPI","Real-time News·Trends","Live Updates",BLUE),
    ("Adzuna/\nJsearch","Jobs·Salary·Market","Aggregated",GREEN),
    ("Paisa\nAPI","Finance·MutualFunds","Read Only",CYAN),
]
for i,(nm,sub,tag,col) in enumerate(apis):
    ax=Inches(0.45+i*2.16); ay=Inches(1.42)
    RR(s,ax,ay,Inches(2.0),Inches(0.72),
       fill=RGBColor(col.red//7,col.green//7,col.blue//7),line=col,lw=Pt(1),r=0.05)
    T(s,nm,ax+Inches(0.08),ay+Inches(0.04),Inches(1.1),Inches(0.42),sz=9,bold=True,color=col)
    T(s,sub,ax+Inches(0.08),ay+Inches(0.42),Inches(1.8),Inches(0.22),sz=7,color=MTXT)
    TAG(s,tag,ax+Inches(1.05),ay+Inches(0.04),
        RGBColor(col.red//6,col.green//6,col.blue//6),col,sz=6.5,h=Pt(14))

# ── 4 Main Sections ──────────────────────────────────────────────────────────
sections = [
    ("1. DATA INPUTS","From user & connected sources",BLUE,[
        ("Health Data","Sleep, Stress, Activity, HRV",GREEN),
        ("Finance Data","Income, Expenses, Savings",AMBER),
        ("Career Data","Role, Skills, Goals",BLUE),
        ("Wearables","Steps, HR, Calories, Sleep",PINK),
        ("Documents","Resumes, Reports, PDFs",MTXT),
    ]),
    ("2. BEYONDSELF\nCORE ENGINE","Deterministic · Explainable · Reliable",PURPLE,[
        ("9 AI Engines","Health·Finance·Career\nSleep·Burnout·Anomaly",PURPLE),
        ("Cross-Domain\nCascades","Inputs → Cascade Engine\n→ Life Score 0–100",CYAN),
    ]),
    ("3. BACKEND\nSERVICES","Spring Boot · REST · Secure",GREEN,[
        ("API Gateway","Routing·Rate Limiting·Auth",BLUE),
        ("Core Services","User·Score·Timeline·Coach",GREEN),
        ("Data Services","CRUD·Sync·Cache·Search",TEAL),
        ("File Upload","CSV·PDF·Images",AMBER),
        ("Notification","Reminders·Alerts·Digest",ORANGE),
    ]),
    ("4. EXPERIENCE\nLAYER","Real insights. Real impact.",ORANGE,[
        ("Life Score","Single score across domains",PURPLE),
        ("Timeline","Track progress over time",CYAN),
        ("AI Insights","Personalised, actionable",GREEN),
        ("Simulator","What-if scenarios",AMBER),
        ("Coach","AI guidance & nudges",ORANGE),
    ]),
]

for i,(sec_title,sec_sub,sec_col,items) in enumerate(sections):
    sx=Inches(0.3+i*3.27); sy=Inches(2.35)
    sw=Inches(3.1); sh=Inches(3.4)
    sc_bg=RGBColor(max(0,sec_col.red//6),max(0,sec_col.green//6),max(0,sec_col.blue//6))
    RR(s,sx,sy,sw,sh, fill=sc_bg,line=sec_col,lw=Pt(1.5),r=0.04)
    BAR(s,sx,sy,sw,sec_col,h=Pt(2.5))
    T(s,sec_title,sx+Inches(0.12),sy+Inches(0.1),sw-Inches(0.24),Inches(0.52),
      sz=11,bold=True,color=sec_col)
    T(s,sec_sub,sx+Inches(0.12),sy+Inches(0.56),sw-Inches(0.24),Inches(0.24),
      sz=8,color=MTXT)
    BAR(s,sx+Inches(0.12),sy+Inches(0.82),sw-Inches(0.24),DIV)
    for j,(it_name,it_sub,it_col) in enumerate(items):
        iy=sy+Inches(0.95+j*0.5)
        ib_bg=RGBColor(it_col.red//7,it_col.green//7,it_col.blue//7)
        RR(s,sx+Inches(0.12),iy,sw-Inches(0.24),Inches(0.42),
           fill=ib_bg,line=it_col,lw=Pt(0.8),r=0.04)
        T(s,it_name,sx+Inches(0.22),iy+Inches(0.04),Inches(1.1),Inches(0.35),
          sz=9,bold=True,color=it_col)
        T(s,it_sub,sx+Inches(1.3),iy+Inches(0.04),sw-Inches(1.55),Inches(0.35),
          sz=8,color=MTXT)

# Security panel
R(s,Inches(13.05),Inches(2.35),Inches(0.25),Inches(3.4),fill=DIV)
sx2=Inches(13.05)
T(s,"🔐 SECURITY", sx2-Inches(2.2),Inches(2.38),Inches(2.1),Inches(0.3),
  sz=9,bold=True,color=RED)
for j,item in enumerate(["JWT Auth","BCrypt","CORS Config","Input Validation","Rate Limiting","PII Protection"]):
    T(s,f"✓ {item}", sx2-Inches(2.2),Inches(2.72+j*0.46),Inches(2.1),Inches(0.35),
      sz=8.5,color=MTXT)

# ── Infrastructure bar bottom ─────────────────────────────────────────────────
R(s,Inches(0.3),Inches(5.85),W-Inches(0.6),Inches(1.35),
  fill=RGBColor(0x12,0x10,0x2A),line=AMBER,lw=Pt(1))
BAR(s,Inches(0.3),Inches(5.85),W-Inches(0.6),AMBER,h=Pt(2.5))
T(s,"5. DATA & INFRASTRUCTURE LAYER  ·  Offline-first · Secure · Scalable",
  Inches(0.45),Inches(5.9),Inches(8),Inches(0.28),sz=9,bold=True,color=AMBER)
T(s,"PostgreSQL (Port 5432)  ·  users · health_records · finance_records · career_records · goals · transactions",
  Inches(0.45),Inches(6.22),Inches(6.2),Inches(0.28),sz=8.5,color=MTXT)
T(s,"LocalStorage · Schema v2+ auto-migration · storageAdapter.js · cross-tab sync",
  Inches(0.45),Inches(6.55),Inches(6.2),Inches(0.28),sz=8.5,color=MTXT)
T(s,"Offline First\nWorks anywhere", Inches(6.8),Inches(5.9),Inches(1.5),Inches(0.6),
  sz=8,bold=True,color=GREEN,align=PP_ALIGN.CENTER)
T(s,"Deterministic AI\nConsistent every time", Inches(8.4),Inches(5.9),Inches(1.7),Inches(0.6),
  sz=8,bold=True,color=PURPLE,align=PP_ALIGN.CENTER)
T(s,"Privacy By Design\nYour data, your control", Inches(10.2),Inches(5.9),Inches(1.7),Inches(0.6),
  sz=8,bold=True,color=CYAN,align=PP_ALIGN.CENTER)
T(s,"Demo Proof\nNo internet, still works", Inches(12.0),Inches(5.9),Inches(1.2),Inches(0.6),
  sz=8,bold=True,color=AMBER,align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — NOVELTY & INNOVATION  (matches screenshot 6)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

T(s,"NOVELTY & INNOVATION  ·  CHALLENGES & BREAKTHROUGHS",
  Inches(0.45),Inches(0.18),W-Inches(4),Inches(0.28),sz=8,bold=True,
  color=RGBColor(0x9C,0x72,0xFF))

# Dual headline
T(s,"What Hit Us.", Inches(0.45),Inches(0.45),Inches(5.5),Inches(0.75),
  sz=32,bold=True,color=ORANGE)
T(s,"  How We Broke Through.", Inches(5.5),Inches(0.45),Inches(7),Inches(0.75),
  sz=32,bold=True,color=PURPL)

brand_pill = Inches(2.7)
RR(s,W-brand_pill-Inches(0.3),Inches(0.15),brand_pill,Inches(0.38),
   fill=RGBColor(0x10,0x0E,0x28),line=GREEN,lw=Pt(1),r=0.4)
T(s,"● ALL 5 BREAKTHROUGHS SHIPPED",
  W-brand_pill-Inches(0.2),Inches(0.18),brand_pill-Inches(0.2),Inches(0.3),
  sz=7.5,bold=True,color=GREEN,align=PP_ALIGN.CENTER)

# Column headers
T(s,"⚡  WHAT HIT US",
  Inches(0.45),Inches(1.22),Inches(5.9),Inches(0.25),sz=8,bold=True,color=AMBER)
BAR(s,Inches(0.45),Inches(1.48),Inches(5.9),AMBER,h=Pt(1.5))
T(s,"✦  HOW WE BROKE THROUGH",
  Inches(7.0),Inches(1.22),Inches(5.9),Inches(0.25),sz=8,bold=True,color=PURPLE)
BAR(s,Inches(7.0),Inches(1.48),Inches(5.9),PURPLE,h=Pt(1.5))

pairs = [
    ("🔀","Data Fragmentation",
     "Fitbit · bank slips · lab reports · LinkedIn — all disconnected silos.",
     "🔗","Unified Feed Parser",
     "Voice + PDF + Photo → 1 Digital Twin. ocrService · resumeService · voiceLogService",
     AMBER, GREEN),
    ("🕸️","Cross-Domain Blindness",
     "Health, Finance, Career in isolation. No app modelled cause-effect chains.",
     "🗺️","Cascade Map Engine",
     "Real cross-domain simulation. Salary↑→Sleep↓→Health↓→Output↓. simulatorEngine.js",
     ORANGE, BLUE),
    ("🎙️","Smart Input Classification",
     "'₹450 on medicine' — Health? Finance? Both? All tagging was manual.",
     "🤖","Universal AI Mic + Parser",
     "Groq llama-3.3-70b detects domain, action & amount. Zero manual tags.",
     PURPLE, TEAL),
    ("🔓","Privacy Exposure",
     "Raw sensitive data — names, accounts, phones — sent directly to AI models.",
     "🛡️","PII Shield",
     "stripPII() fires before every AI call. Identity never reaches the model.",
     RED, CYAN),
    ("❓","Prediction Uncertainty",
     "12-month forecasts shown with same confidence as 1-month ones. System lied.",
     "🧠","Neural Core",
     "Confidence decay: 1mo 92% · 3mo 85% · 6mo 73% · 12mo 48%. Honest AI.",
     ORANGE, PURPL),
]

for j,(pi,pt,pd,si,st,sd,pc,sc) in enumerate(pairs):
    tp = Inches(1.58 + j*1.12)
    rh = Inches(1.0)
    # Problem
    pbg=RGBColor(max(0,pc.red//6),max(0,pc.green//6),max(0,pc.blue//6))
    RR(s,Inches(0.45),tp,Inches(5.9),rh,fill=pbg,line=pc,lw=Pt(1),r=0.04)
    BAR(s,Inches(0.45),tp,Inches(5.9),pc,h=Pt(1.5))
    T(s,f"0{j+1}",Inches(0.52),tp+Inches(0.08),Inches(0.4),Inches(0.55),
      sz=20,bold=True,color=RGBColor(pc.red//4,pc.green//4,pc.blue//4))
    RR(s,Inches(0.98),tp+Inches(0.12),Inches(0.44),Inches(0.44),
       fill=RGBColor(pc.red//5,pc.green//5,pc.blue//5),line=pc,lw=Pt(1),r=0.2)
    T(s,pi,Inches(0.98),tp+Inches(0.1),Inches(0.44),Inches(0.44),sz=18,align=PP_ALIGN.CENTER)
    T(s,pt,Inches(1.5),tp+Inches(0.1),Inches(3.6),Inches(0.32),sz=12,bold=True,color=pc)
    T(s,pd,Inches(1.5),tp+Inches(0.44),Inches(4.6),Inches(0.45),sz=9,color=MTXT,wrap=True)
    # Arrow
    T(s,"▶",Inches(6.4),tp+Inches(0.3),Inches(0.55),Inches(0.4),
      sz=14,color=DTXT,align=PP_ALIGN.CENTER)
    T(s,"SOLVED BY",Inches(6.25),tp+Inches(0.62),Inches(0.8),Inches(0.25),
      sz=6,bold=True,color=DTXT,align=PP_ALIGN.CENTER)
    # Solution
    sbg=RGBColor(max(0,sc.red//6),max(0,sc.green//6),max(0,sc.blue//6))
    RR(s,Inches(7.0),tp,Inches(5.9),rh,fill=sbg,line=sc,lw=Pt(1),r=0.04)
    BAR(s,Inches(7.0),tp,Inches(5.9),sc,h=Pt(1.5))
    RR(s,Inches(7.08),tp+Inches(0.12),Inches(0.44),Inches(0.44),
       fill=RGBColor(sc.red//5,sc.green//5,sc.blue//5),line=sc,lw=Pt(1),r=0.2)
    T(s,si,Inches(7.08),tp+Inches(0.1),Inches(0.44),Inches(0.44),sz=18,align=PP_ALIGN.CENTER)
    T(s,st,Inches(7.6),tp+Inches(0.1),Inches(4.9),Inches(0.32),sz=12,bold=True,color=sc)
    T(s,sd,Inches(7.6),tp+Inches(0.44),Inches(5.0),Inches(0.45),sz=9,color=MTXT,wrap=True)

# Bottom strip
R(s,Inches(0.45),Inches(7.1),W-Inches(0.9),Inches(0.32),
  fill=RGBColor(0x10,0x0E,0x28))
T(s,"💡 Building a Personal Digital Twin isn't about tracking data — it's about connecting every life decision into one intelligent system.",
  Inches(0.55),Inches(7.12),Inches(7.5),Inches(0.26),sz=8.5,color=MTXT)
for k,(n,l_) in enumerate([("5","Breakthroughs"),("0","PII to AI"),("3","Domains"),("0","Manual Tags")]):
    T(s,n, Inches(8.3+k*1.22),Inches(7.07),Inches(1.0),Inches(0.2),
      sz=14,bold=True,color=PURPLE,align=PP_ALIGN.CENTER)
    T(s,l_,Inches(8.1+k*1.22),Inches(7.24),Inches(1.2),Inches(0.14),
      sz=6.5,color=DTXT,align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — TEAM
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)
slide_header(s,"THE TEAM  ·  BUILT BY",
             "Meet the Builders.", PURPLE,
             "BeyondSelf — Personal Digital Twin  ·  4 June 2026")

team_members = [
    ("KB","Khushi Bansal",   "Full-Stack & AI Lead",["React","Groq AI","Python","Architecture"],     PURPLE),
    ("PG","Pavani Gubba",    "Backend & Data Lead",  ["Spring Boot","Score Engines","PostgreSQL","APIs"],CYAN),
    ("SP","Samridhi Pandey", "Design & Frontend Lead",["UI/UX","Tailwind","Framer Motion","Features"],PINK),
]
for i,(init,name,role,skills,col) in enumerate(team_members):
    ml = Inches(1.0 + i*4.1)
    # Avatar
    bg_av=RGBColor(max(0,col.red//5),max(0,col.green//5),max(0,col.blue//5))
    OV(s,ml+Inches(0.65),Inches(1.65),Inches(2.1),Inches(2.1),fill=bg_av,line=col,lw=Pt(2.5))
    # Outer ring
    OV(s,ml+Inches(0.45),Inches(1.45),Inches(2.5),Inches(2.5),
       fill=None,line=RGBColor(col.red//3,col.green//3,col.blue//3),lw=Pt(1.5))
    T(s,init, ml+Inches(0.65),Inches(1.8),Inches(2.1),Inches(1.0),
      sz=38,bold=True,color=col,align=PP_ALIGN.CENTER)
    T(s,name, ml,Inches(3.9),Inches(3.4),Inches(0.5),
      sz=17,bold=True,color=WHITE,align=PP_ALIGN.CENTER)
    T(s,role, ml,Inches(4.42),Inches(3.4),Inches(0.35),
      sz=11,color=col,align=PP_ALIGN.CENTER,italic=True)
    # Skills
    for k,sk in enumerate(skills):
        sk_l=ml+Inches(0.08+k%2*1.65); sk_t=Inches(4.9+(k//2)*0.52)
        sk_bg=RGBColor(col.red//5,col.green//5,col.blue//5)
        RR(s,sk_l,sk_t,Inches(1.52),Inches(0.38),fill=sk_bg,line=col,lw=Pt(0.8),r=0.3)
        T(s,sk, sk_l,sk_t,Inches(1.52),Inches(0.38),
          sz=9,bold=True,color=col,align=PP_ALIGN.CENTER)

T(s,"Built with:  React 19 · Vite 8 · Spring Boot 3.2 · Groq AI · Gemini 2.0 · Tailwind v4 · Tesseract.js · Nutritionix · GitHub API",
  Inches(0.5),Inches(6.8),W-Inches(1),Inches(0.3),
  sz=9,color=DTXT,align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CLOSING
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

R(s,0,0,W,Inches(0.05),fill=PURPLE)  # top bar

T(s,"Building the world's first",
  Inches(0.5),Inches(0.9),W-Inches(1),Inches(0.45),
  sz=18,color=MTXT,align=PP_ALIGN.CENTER)
T(s,"Personal Digital Twin",
  Inches(0.5),Inches(1.35),W-Inches(1),Inches(1.2),
  sz=54,bold=True,color=PURPLE,align=PP_ALIGN.CENTER)
T(s,"that connects every life decision into one intelligent system.",
  Inches(1),Inches(2.6),W-Inches(2),Inches(0.5),
  sz=16,color=MTXT,align=PP_ALIGN.CENTER,italic=True)

BAR(s,Inches(4.5),Inches(3.28),Inches(4.33),PURPLE,h=Pt(2))

# Stats row
for i,(n,l_) in enumerate([("9","Original\nFeatures"),("3","Domains\nUnified"),
                             ("0","PII to AI"),("847","Simulations\nRun"),("∞","Cascade\nPaths")]):
    sl=Inches(0.8+i*2.55)
    RR(s,sl,Inches(3.5),Inches(2.2),Inches(1.55),
       fill=RGBColor(0x14,0x11,0x30),line=DIV,lw=Pt(1),r=0.08)
    T(s,n,  sl,Inches(3.56),Inches(2.2),Inches(0.8),
      sz=38,bold=True,color=PURPLE,align=PP_ALIGN.CENTER)
    T(s,l_, sl,Inches(4.38),Inches(2.2),Inches(0.55),
      sz=9,color=MTXT,align=PP_ALIGN.CENTER)

T(s,"BeyondSelf  ·  4 June 2026  ·  Khushi Bansal  ·  Pavani Gubba  ·  Samridhi Pandey",
  Inches(0.5),Inches(6.8),W-Inches(1),Inches(0.3),
  sz=10,color=DTXT,align=PP_ALIGN.CENTER)
R(s,0,H-Inches(0.05),W,Inches(0.05),fill=PURPLE)

# ─── Save ────────────────────────────────────────────────────────────────────
out = r"C:\Users\Saanvi Jaiswal\Downloads\beyondselff\BeyondSelf_Pitch_Deck.pptx"
prs.save(out)
print(f"✓ Saved: {out}")
print(f"✓ Slides: {len(prs.slides)}")
