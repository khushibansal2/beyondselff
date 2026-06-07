"""
BeyondSelf — Screenshot every HTML slide → embed into PPTX
Pixel-perfect: exactly what Chrome renders.
"""
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches
import os, sys

BASE = r'C:\Users\Saanvi Jaiswal\Downloads\beyondselff'
OUT  = os.path.join(BASE, 'BeyondSelf_Full_Deck.pptx')
IMGS = os.path.join(BASE, '_slide_imgs')
os.makedirs(IMGS, exist_ok=True)

# ── Ordered slide list (main deck first, then supplementary) ─────────────────
FILES = [
    'agenda-slide.html',
    'slide-01-problem.html',
    'slide-02-architecture.html',
    'slide-03-uiux.html',
    'slide-04-dataflow.html',
    'slide-05-techstack.html',
    'tech-stack-slide.html',
    'slide-06-outcome.html',
    'slide-07-novelty.html',
    'slide-08-challenges.html',
    'slide-09-team.html',
    'slide-10-references.html',
    'dfd_2slides.html',
    'dfd_presentation.html',
    'novelty_innovation.html',
    'challenges_breakthroughs.html',
    'novelty_features.html',
    'BeyondSelf_Architecture.html',
]

W, H = 1400, 900
screenshots = []   # ordered list of PNG paths

print("=== Screenshotting HTML slides ===")
with sync_playwright() as p:
    browser = p.chromium.launch(headless=True)

    for fname in FILES:
        fpath = os.path.join(BASE, fname)
        if not os.path.exists(fpath):
            print(f'  SKIP (not found): {fname}')
            continue

        # file:// URL — forward slashes required
        url = 'file:///' + fpath.replace('\\', '/')
        page = browser.new_page(viewport={'width': W, 'height': H})
        try:
            page.goto(url, wait_until='networkidle', timeout=20000)
        except Exception:
            page.goto(url, timeout=20000)
        page.wait_for_timeout(1200)  # fonts + animations settle

        # Count how many .slide elements this file has
        n = page.evaluate(
            "document.querySelectorAll('.slide').length || 1"
        )

        if n > 1:
            for i in range(n):
                # Activate slide i via DOM manipulation
                page.evaluate(f"""
                    const slides = document.querySelectorAll('.slide');
                    slides.forEach(s => {{
                        s.classList.remove('active');
                        s.style.display = 'none';
                    }});
                    slides[{i}].classList.add('active');
                    slides[{i}].style.display = 'flex';
                    const ctr = document.getElementById('ctr');
                    if (ctr) ctr.textContent = '{i+1} / ' + slides.length;
                    const nav = document.querySelector('.nav');
                    if (nav) nav.style.display = 'none';
                """)
                page.wait_for_timeout(300)
                img = os.path.join(IMGS, f'{fname[:-5]}_p{i+1}.png')
                page.screenshot(path=img,
                                clip={'x': 0, 'y': 0, 'width': W, 'height': H})
                screenshots.append(img)
                print(f'  [{len(screenshots):02d}] {fname}  ->  page {i+1}/{n}')
        else:
            # Hide nav bar if present so it doesn't overlap the slide
            page.evaluate("""
                const nav = document.querySelector('.nav');
                if (nav) nav.style.display = 'none';
            """)
            img = os.path.join(IMGS, f'{fname[:-5]}.png')
            page.screenshot(path=img,
                            clip={'x': 0, 'y': 0, 'width': W, 'height': H})
            screenshots.append(img)
            print(f'  [{len(screenshots):02d}] {fname}')

        page.close()

    browser.close()

print(f'\n=== Building PPTX with {len(screenshots)} slides ===')
prs = Presentation()
prs.slide_width  = Inches(13.33)   # 16:9 widescreen
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

for img_path in screenshots:
    if not os.path.exists(img_path):
        print(f'  MISSING image: {img_path}')
        continue
    sl = prs.slides.add_slide(BLANK)
    # Fill entire slide with the screenshot
    sl.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)

prs.save(OUT)
print(f'\nSaved: {OUT}')
print(f'Total slides: {len(prs.slides)}')
